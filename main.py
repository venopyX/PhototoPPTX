#TOKEN = '6027255433:AAEOSUy7MSQyF1zThWHnkXUkTnjWpu0Wxsw'


import os
import requests
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import telebot

# Replace with your Telegram bot token
TOKEN = '6027255433:AAEOSUy7MSQyF1zThWHnkXUkTnjWpu0Wxsw'

# Create a new Telebot instance
bot = telebot.TeleBot(TOKEN)

# Define a function to handle the /start command
@bot.message_handler(commands=['start'])
def handle_start(message):
    bot.reply_to(message, 'Welcome to the PowerPoint bot! Please send me an image and I will create a PowerPoint presentation with it as the background.')

# Define a function to handle incoming messages
@bot.message_handler(content_types=['photo'])
def handle_message(message):
    chat_id = message.chat.id
    file_id = message.photo[-1].file_id
    # Use the Telegram Bot API to download the photo
    file_info = bot.get_file(file_id)
    file_url = f'https://api.telegram.org/file/bot{TOKEN}/{file_info.file_path}'
    response = requests.get(file_url)
    # Save the photo to a local directory
    with open('image.jpg', 'wb') as f:
        f.write(response.content)
    # Process the image and create a PowerPoint presentation
    try:
        image = Image.open('image.jpg')
        width, height = image.size
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = (255, 255, 255)
        left = top = 0
        slide.shapes.add_picture('image.jpg', left, top, width=Inches(10), height=Inches(7.5))
        # Add a text box to the slide
        textbox = slide.shapes.add_textbox(Inches(0), Inches(6.5), Inches(10), Inches(0.5))
        textbox.text = 'Created with: GCPPT_bot'
        # Save the PowerPoint presentation to a file
        prs.save('presentation.pptx')
        # Send the PowerPoint presentation back to the user
        with open('presentation.pptx', 'rb') as f:
            bot.send_document(chat_id, f)
    except Exception as e:
        bot.reply_to(message, f'Error: {e}')

# Define a function to handle errors
@bot.message_handler(func=lambda message: True)
def handle_error(message):
    bot.reply_to(message, 'Sorry, I could not process your request. Please send me an image and I will create a PowerPoint presentation with it as the background.')

# Start the bot
bot.polling()
